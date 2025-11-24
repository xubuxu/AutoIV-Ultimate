"""
AutoIV-Ultimate Reporting Module - PowerPoint Builder

Creates comprehensive PowerPoint presentations with images and data tables.
"""
import logging
from pathlib import Path
from typing import Dict, Optional

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)


class PowerPointBuilder:
    """Generates PowerPoint presentations with images and tables."""
    
    def __init__(self, output_dir: str = "output"):
        """
        Initialize PowerPointBuilder.
        
        Args:
            output_dir: Directory to save PowerPoint files
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        logger.info("PowerPointBuilder initialized")
    
    def create_report(self, stats_df: pd.DataFrame, champion_df: pd.DataFrame,
                     img_paths: Dict[str, str], filename: str = "IV_Analysis_Slides.pptx") -> str:
        """
        Create comprehensive PowerPoint presentation.
        
        Args:
            stats_df: Statistics dataframe
            champion_df: Champion cells dataframe
            img_paths: Dictionary mapping plot names to file paths
            filename: Output filename
            
        Returns:
            Path to saved PowerPoint file
        """
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            self._add_title_slide(prs, "IV Analysis Report")
            
            # Statistics slide
            self._add_table_slide(prs, "Statistical Summary", stats_df)
            
            # Champion cells slide
            self._add_table_slide(prs, "Champion Cells", champion_df)
            
            # Plot slides
            self._add_plot_slides(prs, img_paths)
            
            # Save
            output_path = self.output_dir / filename
            prs.save(output_path)
            
            logger.info(f"PowerPoint report saved: {output_path}")
            return str(output_path)
            
        except Exception as e:
            logger.error(f"PowerPoint report generation failed: {e}")
            return ""
    
    def _add_title_slide(self, prs: Presentation, title: str):
        """Add title slide."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = "Automated IV Characterization Analysis"
    
    def _add_table_slide(self, prs: Presentation, title: str, df: pd.DataFrame):
        """Add slide with data table."""
        if df.empty:
            return
        
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        
        # Limit columns and rows to fit on slide
        df_display = df.head(10).iloc[:, :8]  # Max 10 rows, 8 columns
        
        # Add table
        rows, cols = df_display.shape[0] + 1, df_display.shape[1]
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Add headers
        for i, column in enumerate(df_display.columns):
            cell = table.cell(0, i)
            cell.text = str(column)
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].font.bold = True
        
        # Add data
        for row_idx, row in enumerate(df_display.itertuples(index=False), start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
                cell.text_frame.paragraphs[0].font.size = Pt(9)
    
    def _add_plot_slides(self, prs: Presentation, img_paths: Dict[str, str]):
        """Add slides with plot images."""
        plot_titles = {
            'boxplot': 'Parameter Boxplots',
            'histogram': 'Efficiency Distribution',
            'trend': 'Batch Trends',
            'yield': 'Yield Distribution',
            'jv_curves': 'Champion J-V Curves',
            'voc_jsc': 'Voc vs Jsc Correlation',
            'drivers': 'Efficiency Drivers',
            'resistance': 'Resistance Analysis',
            'corr_matrix': 'Correlation Matrix'
        }
        
        for key, path in img_paths.items():
            if path and Path(path).exists():
                title = plot_titles.get(key, key.replace('_', ' ').title())
                
                slide_layout = prs.slide_layouts[5]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
                title_frame = title_box.text_frame
                title_frame.text = title
                title_frame.paragraphs[0].font.size = Pt(24)
                title_frame.paragraphs[0].font.bold = True
                title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Add image
                try:
                    left = Inches(1)
                    top = Inches(1.2)
                    height = Inches(5.5)
                    slide.shapes.add_picture(path, left, top, height=height)
                except Exception as e:
                    logger.warning(f"Failed to add image {path}: {e}")
